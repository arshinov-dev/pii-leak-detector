"""
Universal OCR module.

Features:
- Images:   TIF, TIFF, JPEG, JPG, PNG, GIF
- PDFs:     text + scanned pages
- Videos:   MP4 (sample frames, parallel OCR)
- Slides:   PPT / PPTX (text + images via OCR)

Engine: Tesseract (pytesseract).
Languages: Russian, English, Czech (OCR_LANG).

from ocr_extractor import extract_text

#texts = extract_text("ckp.pdf")      # → ['стр 1...', 'стр 2...']
texts = extract_text("CS10.mp4")     # → ['срез 1...', 'срез 2...']
"""

import io
import re
import subprocess
import tempfile
from collections import Counter
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Union, Iterable, List

import cv2
import fitz  # PyMuPDF
import numpy as np
import pytesseract
from PIL import Image, ImageFilter, ImageEnhance
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Global configuration
LIBREOFFICE_PATH: str = "soffice"

OCR_LANG: str = "rus+eng+ces"
PDF_RENDER_DPI: int = 150

PSM_DOC: int = 3      # page layout analysis
PSM_VIDEO: int = 6    # single block of text

MIN_LINE_LENGTH: int = 3
UPSCALE_TARGET_PX: int = 2000

VIDEO_SAMPLE_INTERVAL_SEC: float = 1.0
VIDEO_MAX_WORKERS: int = 4
VIDEO_MIN_ALPHA_RATIO: float = 0.45

MIN_RUN_HEADER_SHARE: float = 0.30


def _clean_text(text: str) -> str:
    """
    Normalize OCR text:
    - remove control chars and private-use Unicode
    - drop soft hyphens and zero-width chars
    - collapse newlines to spaces and normalize spaces.
    """
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
    text = re.sub(r'[\uE000-\uF8FF]', '', text)
    text = text.replace('\xad', '').replace('\u200b', '').replace('\u200c', '').replace('\u200d', '')
    text = re.sub(r'[\xa0\u2002\u2003\u2009\u202f]', ' ', text)
    text = re.sub(r'-\n\s*', '', text)
    text = re.sub(r'\n+', ' ', text)
    text = re.sub(r'[ \t]+', ' ', text)
    return text.strip()


def _remove_running_headers(pages: List[str]) -> List[str]:
    """
    Remove lines that repeat on many PDF pages (headers/footers).
    Any line that appears on >= MIN_RUN_HEADER_SHARE pages is treated as noise.
    """
    total = len(pages)
    if total < 3:
        return pages

    line_page_count: Counter[str] = Counter()
    for page in pages:
        for line in set(page.splitlines()):
            s = line.strip()
            if s:
                line_page_count[s] += 1

    threshold = max(3, int(total * MIN_RUN_HEADER_SHARE))
    headers = {line for line, cnt in line_page_count.items() if cnt >= threshold}

    cleaned: List[str] = []
    for page in pages:
        lines = [l for l in page.splitlines() if l.strip() not in headers]
        cleaned.append(" ".join(lines).strip())
    return cleaned


def _score_text(text: str) -> float:
    """
    Heuristic text quality score in [0, 1]:
    - fraction of alphanumeric/space characters
    - bonus for words of length >= 3.
    """
    if not text:
        return 0.0
    total = len(text)
    alpha = sum(1 for c in text if c.isalnum() or c.isspace())
    ratio = alpha / total
    words = text.split()
    long_words = sum(1 for w in words if len(w) >= 3)
    word_bonus = min(long_words / max(len(words), 1), 1.0) * 0.3
    return ratio * 0.7 + word_bonus


def _is_real_text(text: str) -> bool:
    """
    Check that a line looks like real text, not noise:
    - sufficient share of alphanumeric characters
    - has words of length >= 3
    - contains at least one letter (not only digits/punctuation).
    """
    if not text or len(text) < 4:
        return False
    alpha = sum(1 for c in text if c.isalnum() or c.isspace())
    if alpha / len(text) < VIDEO_MIN_ALPHA_RATIO:
        return False
    words = [w for w in text.split() if len(w) >= 3]
    if not words:
        return False
    if not re.search(r'[a-zA-Zа-яА-ЯёЁčšžřůýáéíóúěďťň]', text):
        return False
    return True


def _deduplicate(blocks: Iterable[str]) -> List[str]:
    """
    Remove consecutive duplicate blocks.
    Keeps first occurrence in each run.
    """
    result: List[str] = []
    prev: str | None = None
    for b in blocks:
        if b and b != prev:
            result.append(b)
            prev = b
    return result


def _upscale(arr: np.ndarray) -> np.ndarray:
    """
    Upscale array so that the longest side equals UPSCALE_TARGET_PX.
    Uses bicubic interpolation; returns new array.
    """
    h, w = arr.shape[:2]
    if max(h, w) < UPSCALE_TARGET_PX:
        scale = UPSCALE_TARGET_PX / max(h, w)
        arr = cv2.resize(arr, (int(w * scale), int(h * scale)), interpolation=cv2.INTER_CUBIC)
    return arr


def _to_gray(arr: np.ndarray) -> np.ndarray:
    """Convert RGB array to grayscale if needed."""
    return cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY) if arr.ndim == 3 else arr


def _prep_otsu(arr: np.ndarray) -> Image.Image:
    """
    Generic preprocessing:
    denoise → local contrast (CLAHE) → global Otsu binarization.
    Works well for most printed text.
    """
    arr = _upscale(arr)
    gray = _to_gray(arr)
    gray = cv2.fastNlMeansDenoising(gray, h=10)
    clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
    gray = clahe.apply(gray)
    _, bin_img = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    bin_img = cv2.copyMakeBorder(
        bin_img, 20, 20, 20, 20, cv2.BORDER_CONSTANT, value=255
    )
    return Image.fromarray(bin_img)


def _prep_adaptive(arr: np.ndarray) -> Image.Image:
    """
    Adaptive binarization for uneven lighting.
    Use when Otsu fails due to gradients or strong shadows.
    """
    arr = _upscale(arr)
    gray = _to_gray(arr)
    gray = cv2.fastNlMeansDenoising(gray, h=12)
    bin_img = cv2.adaptiveThreshold(
        gray,
        255,
        cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
        cv2.THRESH_BINARY,
        blockSize=31,
        C=11,
    )
    bin_img = cv2.copyMakeBorder(
        bin_img, 20, 20, 20, 20, cv2.BORDER_CONSTANT, value=255
    )
    return Image.fromarray(bin_img)


def _prep_document(arr: np.ndarray) -> Image.Image:
    """
    Document‑specific preprocessing (e.g., passport/ID):
    estimate background via dilation, subtract it, enhance contrast,
    then apply Otsu binarization.
    """
    arr = _upscale(arr)
    gray = _to_gray(arr)
    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (25, 25))
    bg = cv2.morphologyEx(gray, cv2.MORPH_DILATE, kernel)
    diff = cv2.absdiff(gray, bg)
    diff = 255 - diff
    clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
    enhanced = clahe.apply(diff)
    _, bin_img = cv2.threshold(
        enhanced, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU
    )
    bin_img = cv2.copyMakeBorder(
        bin_img, 20, 20, 20, 20, cv2.BORDER_CONSTANT, value=255
    )
    return Image.fromarray(bin_img)


def _prep_simple(arr: np.ndarray) -> Image.Image:
    """
    Simple fallback preprocessing:
    grayscale → sharpen → slight contrast boost.
    """
    pil = Image.fromarray(arr).convert("L")
    pil = pil.filter(ImageFilter.SHARPEN)
    pil = ImageEnhance.Contrast(pil).enhance(1.4)
    return pil


PREP_FUNCS = [_prep_otsu, _prep_adaptive, _prep_document, _prep_simple]


def _detect_and_crop_card(arr: np.ndarray) -> np.ndarray:
    """
    Try to detect a rectangular ID card / document near the bottom of the frame
    and rectify it via perspective transform.

    If detection fails, returns the original array.
    """
    h, w = arr.shape[:2]
    gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
    blur = cv2.GaussianBlur(gray, (7, 7), 0)

    edges = cv2.Canny(blur, 30, 90)
    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (5, 5))
    edges = cv2.dilate(edges, kernel, iterations=2)
    edges = cv2.morphologyEx(edges, cv2.MORPH_CLOSE, kernel, iterations=2)

    mask = np.zeros_like(edges)
    mask[int(h * 0.35):, :] = 255
    edges = cv2.bitwise_and(edges, mask)

    contours, _ = cv2.findContours(
        edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE
    )
    if not contours:
        return arr

    contours = sorted(contours, key=cv2.contourArea, reverse=True)

    card_contour: np.ndarray | None = None
    for cnt in contours[:15]:
        area = cv2.contourArea(cnt)
        if area < (h * w * 0.05) or area > (h * w * 0.80):
            continue
        peri = cv2.arcLength(cnt, True)
        approx = cv2.approxPolyDP(cnt, 0.015 * peri, True)
        if len(approx) == 4 and cv2.isContourConvex(approx):
            card_contour = approx.reshape(-1, 2).astype(np.float32)
            break

    if card_contour is None:
        rect = cv2.minAreaRect(contours[0])
        box = cv2.boxPoints(rect)
        card_contour = box.astype(np.float32)

    if card_contour.shape[0] != 4:
        return arr

    pts = card_contour
    rect = np.zeros((4, 2), dtype=np.float32)
    s = pts.sum(axis=1)
    diff = np.diff(pts, axis=1)
    rect[0] = pts[np.argmin(s)]      # top-left
    rect[2] = pts[np.argmax(s)]      # bottom-right
    rect[1] = pts[np.argmin(diff)]   # top-right
    rect[3] = pts[np.argmax(diff)]   # bottom-left

    wa = np.linalg.norm(rect[2] - rect[3])
    wb = np.linalg.norm(rect[1] - rect[0])
    ha = np.linalg.norm(rect[1] - rect[2])
    hb = np.linalg.norm(rect[0] - rect[3])
    out_w, out_h = int(max(wa, wb)), int(max(ha, hb))

    if out_w < 150 or out_h < 80:
        return arr

    dst = np.array(
        [[0, 0], [out_w - 1, 0], [out_w - 1, out_h - 1], [0, out_h - 1]],
        dtype=np.float32,
    )
    M = cv2.getPerspectiveTransform(rect, dst)
    warped = cv2.warpPerspective(arr, M, (out_w, out_h))
    return warped


def _run_tesseract(pil_img: Image.Image, psm: int) -> str:
    """
    Call Tesseract on a PIL image with given PSM and shared OCR_LANG.
    Returns cleaned text.
    """
    raw = pytesseract.image_to_string(
        pil_img, lang=OCR_LANG, config=f"--oem 1 --psm {psm}"
    )
    return _clean_text(raw)


def _ocr_image_arr(arr: np.ndarray, psm: int, try_card: bool = True) -> str:
    """
    Perform OCR on an RGB image array:

    - optionally detect and crop an ID card/document first;
    - for each candidate (crop and full frame) run multiple preprocessing
      pipelines and choose the variant with highest score.
    """
    sources: list[np.ndarray] = []

    if try_card:
        card = _detect_and_crop_card(arr)
        sources.append(card)
        if not np.array_equal(card, arr):
            sources.append(arr)
    else:
        sources = [arr]

    best_text = ""
    best_score = 0.0

    for src in sources:
        for prep in PREP_FUNCS:
            try:
                pil_img = prep(src)
                text = _run_tesseract(pil_img, psm=psm)
                score = _score_text(text)
                if score > best_score:
                    best_score, best_text = score, text
            except Exception:
                continue
        if best_score >= 0.7:
            break

    lines = [
        l.strip()
        for l in best_text.splitlines()
        if len(l.strip()) >= MIN_LINE_LENGTH
    ]
    return " ".join(lines)


def _ocr_from_bytes(image_bytes: bytes) -> str:
    """
    OCR for images embedded in PPTX.
    Skips too small images to avoid noise.
    """
    try:
        img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    except Exception:
        return ""
    if img.width < 80 or img.height < 80:
        return ""
    return _ocr_image_arr(np.array(img), psm=PSM_DOC, try_card=False)


def _process_image(path: str) -> List[str]:
    """
    Process image files (TIF / TIFF / JPEG / JPG / PNG / GIF).
    For GIF returns one block per frame.
    """
    img = Image.open(path)
    frames: list[Image.Image] = []

    try:
        while True:
            frames.append(img.copy().convert("RGB"))
            img.seek(img.tell() + 1)
    except EOFError:
        pass
    if not frames:
        frames = [img.convert("RGB")]

    texts: List[str] = []
    for f in frames:
        text = _ocr_image_arr(np.array(f), psm=PSM_DOC, try_card=True)
        if text:
            texts.append(text)
    return texts or [""]


def _process_pdf(path: str) -> List[str]:
    """
    Process PDF:
    - directly extract digital text where available
    - for scanned pages render to image and run OCR.
    """
    doc = fitz.open(path)
    pages_text: List[str] = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        digital = _clean_text(page.get_text("text", clip=page.rect).strip())
        if digital and len(digital) > 20:
            pages_text.append(digital)
            continue

        zoom = PDF_RENDER_DPI / 72
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        pages_text.append(
            _ocr_image_arr(np.array(img), psm=PSM_DOC, try_card=False)
        )

    doc.close()
    return _remove_running_headers(pages_text)


def _process_video(path: str) -> List[str]:
    """
    Process MP4 video:

    - sample frames every VIDEO_SAMPLE_INTERVAL_SEC seconds;
    - for each frame, attempt ID/document detection and OCR;
    - filter obvious noise and deduplicate consecutive blocks.
    """
    cap = cv2.VideoCapture(path)
    if not cap.isOpened():
        raise IOError(f"Cannot open video: {path}")

    fps = cap.get(cv2.CAP_PROP_FPS) or 25.0
    step = max(1, int(fps * VIDEO_SAMPLE_INTERVAL_SEC))

    frames: list[np.ndarray] = []
    idx = 0
    while True:
        ret, frame = cap.read()
        if not ret:
            break
        if idx % step == 0:
            frames.append(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
        idx += 1
    cap.release()

    def process_frame(arr: np.ndarray) -> str:
        text = _ocr_image_arr(arr, psm=PSM_VIDEO, try_card=True)
        good_lines = [l for l in text.splitlines() if _is_real_text(l)]
        result = " ".join(good_lines)
        return result if _is_real_text(result) else ""

    with ThreadPoolExecutor(max_workers=VIDEO_MAX_WORKERS) as pool:
        texts = list(pool.map(process_frame, frames))

    texts = [t for t in texts if t]
    return _deduplicate(texts) or ["(текст в видео не обнаружен)"]


def _iter_shape_images(shape) -> Iterable[bytes]:
    """
    Yield image blobs from a PPTX shape:
    - direct pictures,
    - pictures inside groups,
    - images used as shape fills.
    """
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            yield shape.image.blob
        except Exception:
            return
        return

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from _iter_shape_images(child)
        return

    try:
        fill = shape.fill
        if fill and fill._fill.tag.endswith("}blipFill"):
            blip = fill._fill.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
            )
            if blip is not None:
                r_id = blip.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                )
                part = shape.part.related_parts.get(r_id) if r_id else None
                if part:
                    yield part.blob
    except Exception:
        return


def _extract_pptx_text(path: str) -> List[str]:
    """
    Extract text from PPTX slides:
    - plain text (text frames, tables)
    - OCR for images on slides.
    """
    prs = Presentation(path)
    slides_text: List[str] = []

    for slide in prs.slides:
        parts: List[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = _clean_text(" ".join(run.text for run in para.runs))
                    if len(line) >= MIN_LINE_LENGTH:
                        parts.append(line)

            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        _clean_text(cell.text)
                        for cell in row.cells
                        if cell.text.strip()
                    )
                    if len(row_text) >= MIN_LINE_LENGTH:
                        parts.append(row_text)

            for blob in _iter_shape_images(shape):
                img_text = _ocr_from_bytes(blob)
                if img_text:
                    parts.append(f"[OCR] {img_text}")

        slides_text.append(" ".join(parts) if parts else "")

    return slides_text


def _ppt_to_pdf(path: str, out_dir: str) -> str:
    """
    Convert PPT / PPTX to PDF using LibreOffice in headless mode.
    Returns path to generated PDF in out_dir.
    """
    result = subprocess.run(
        [
            LIBREOFFICE_PATH,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            out_dir,
            path,
        ],
        capture_output=True,
        text=True,
        timeout=120,
    )
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice error:\n{result.stderr}")

    pdf_path = str(Path(out_dir) / f"{Path(path).stem}.pdf")
    if not Path(pdf_path).exists():
        raise FileNotFoundError(f"LibreOffice did not create PDF: {pdf_path}")
    return pdf_path


def _process_presentation(path: str) -> List[str]:
    """
    Process PPTX / PPT:
    - for PPTX: extract text directly + OCR images, fallback via PDF
      if slides mostly empty;
    - for PPT: always convert to PDF and process as document.
    """
    ext = Path(path).suffix.lower()

    if ext == ".ppt":
        with tempfile.TemporaryDirectory() as tmp:
            return _process_pdf(_ppt_to_pdf(path, tmp))

    direct = _extract_pptx_text(path)
    empty_count = sum(1 for t in direct if not t.strip())

    if empty_count > len(direct) / 2:
        with tempfile.TemporaryDirectory() as tmp:
            pdf_text = _process_pdf(_ppt_to_pdf(path, tmp))
        merged: List[str] = []
        for i, d in enumerate(direct):
            merged.append(d or (pdf_text[i] if i < len(pdf_text) else ""))
        return merged

    return direct


def extract_text(file_path: Union[str, Path]) -> List[str]:
    """
    Main public API.

    Parameters
    ----------
    file_path:
        Path to file. Supported extensions:
        - Images:  .tif, .tiff, .jpeg, .jpg, .png, .gif
        - PDFs:    .pdf
        - Video:   .mp4
        - Slides:  .ppt, .pptx

    Returns
    -------
    list of str:
        Text blocks depending on file type:
        - Image:        [full_text] or per frame for GIF.
        - PDF:          [page1_text, page2_text, ...].
        - Video (MP4):  [segment1, segment2, ...] per sampled frame.
        - Presentation: [slide1_text, slide2_text, ...].
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower().lstrip(".")

    if ext in {"tif", "tiff", "jpeg", "jpg", "png", "gif"}:
        return _process_image(str(path))
    if ext == "pdf":
        return _process_pdf(str(path))
    if ext == "mp4":
        return _process_video(str(path))
    if ext in {"pptx", "ppt"}:
        return _process_presentation(str(path))

    raise ValueError(
        f"Unsupported extension .{ext}. "
        f"Supported: tif, tiff, jpeg, jpg, png, gif, mp4, pdf, pptx, ppt"
    )


def _main() -> None:
    """
    Simple CLI entrypoint.

    Usage:
        python ocr_extractor.py <file> [--json]

    Prints recognized text blocks to stdout.
    """
    import sys
    import json

    if len(sys.argv) < 2:
        print("Usage: python ocr_extractor.py <file> [--json]")
        raise SystemExit(1)

    file_name = sys.argv[1]
    print(f"[OCR] Processing: {file_name}")

    try:
        blocks = extract_text(file_name)
        if "--json" in sys.argv:
            print(json.dumps(blocks, ensure_ascii=False, indent=2))
        else:
            for i, text in enumerate(blocks, 1):
                label = f"[Block {i}]" if len(blocks) > 1 else "[Result]"
                print(f"\n{label}\n" + "─" * 40)
                print(text or "(no text detected)")
        print(f"\nDone. Blocks: {len(blocks)}")
    except Exception as exc:
        print(f"[Error] {exc}")
        raise


if __name__ == "__main__":
    _main()